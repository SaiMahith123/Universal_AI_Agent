import streamlit as st
from groq import Groq
import pypdf as PyPDF2
from pptx import Presentation
import base64
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import tempfile
import os
from streamlit_TTS import auto_play, text_to_audio

# --- 1. CONFIGURATION & CUSTOM STYLING ---
st.set_page_config(
    page_title="Mahi's Universal Groq AI", 
    layout="wide", 
    initial_sidebar_state="expanded"
)

# Custom CSS for UI layout and visual cards
st.markdown("""
    <style>
        .stApp { flex-direction: row-reverse; }
        [data-testid="stSidebar"] { 
            left: auto !important; 
            right: 0 !important; 
            border-left: 1px solid rgba(250, 250, 250, 0.1); 
        }
        .stChatInputContainer { position: fixed; bottom: 20px; z-index: 1000; }
        .main .block-container { padding-bottom: 150px; }
        .file-card { 
            padding: 12px; 
            border-radius: 10px; 
            border: 1px solid #4A4A4A; 
            background-color: #1E1E1E; 
            display: flex; 
            align-items: center; 
            gap: 10px; 
            margin-bottom: 10px; 
        }
    </style>
""", unsafe_allow_html=True)

st.title("🤖 Mahi's Universal Groq AI")

# --- 2. API & STATE INITIALIZATION ---
if "GROQ_API_KEY" in st.secrets:
    client = Groq(api_key=st.secrets["GROQ_API_KEY"])
else:
    st.error("❌ GROQ_API_KEY missing in Secrets!")
    st.stop()

# Persistent state to ensure memory and model switching work
if "messages" not in st.session_state:
    st.session_state.messages = []
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "image_list" not in st.session_state:
    st.session_state.image_list = []
if "last_audio" not in st.session_state:
    st.session_state.last_audio = None

@st.cache_resource
def get_embed_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

# --- 3. MULTI-DOCUMENT PROCESSING (PDF & PPT) ---
def process_docs(uploaded_file):
    chunks, metadata = [], []
    # PDF parsing logic
    if uploaded_file.name.lower().endswith('.pdf'):
        reader = PyPDF2.PdfReader(uploaded_file)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                p_chunks = [text[j:j+1000] for j in range(0, len(text), 800)]
                chunks.extend(p_chunks)
                metadata.extend([f"PDF Pg {i+1}"] * len(p_chunks))
    # PPT parsing logic
    elif uploaded_file.name.lower().endswith(('.ppt', '.pptx')):
        prs = Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            text = " ".join([sh.text for sh in slide.shapes if hasattr(sh, "text")])
            if text:
                s_chunks = [text[j:j+1000] for j in range(0, len(text), 800)]
                chunks.extend(s_chunks)
                metadata.extend([f"PPT Slide {i+1}"] * len(s_chunks))
    
    # Vector indexing
    if chunks:
        em = get_embed_model()
        embeddings = em.encode(chunks)
        idx = faiss.IndexFlatL2(embeddings.shape[1])
        idx.add(np.array(embeddings).astype('float32'))
        st.session_state.vector_store = {"index": idx, "chunks": chunks, "sources": metadata}
        return True
    return False

# --- 4. SIDEBAR UI ---
with st.sidebar:
    st.header("⚙️ Settings")
    
    # UPDATED: Verified Active Model IDs for Groq as of March 2026
    model_choice = st.selectbox("Select Brain:", [
        "llama-3.3-70b-versatile", 
        "meta-llama/llama-4-scout-17b-16e-instruct", 
        "deepseek-r1-distill-llama-70b-v2" # Verified replacement
    ])
    
    # Vertical Guide
    st.info("💡 **Model Guide:**\n- **Llama 3.3:** Best for PDFs & PPTs\n- **Llama 4:** Vision/Images\n- **DeepSeek:** Advanced Logic/Math")
    
    voice_on = st.toggle("🔊 Auto-Play AI Voice", value=True)
    st.divider()
    
    # Download Chat History logic
    if st.session_state.messages:
        chat_text = "\n\n".join([f"{m['role'].upper()}: {m['content']}" for m in st.session_state.messages])
        st.download_button("📥 Download Chat History", data=chat_text, file_name="mahi_chat.txt")

    if st.button("🗑️ Clear Chat"):
        st.session_state.messages, st.session_state.vector_store, st.session_state.image_list = [], None, []
        st.rerun()
    st.caption("Developed by T Sai Mahit | B.E in AI-ML (2021-25)")

# --- 5. CHAT DISPLAY ---
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        if "files" in msg:
            for f_info in msg["files"]:
                st.markdown(f'<div class="file-card">📄 <b>{f_info["name"]}</b></div>', unsafe_allow_html=True)
                if f_info['type'] == "Image":
                    st.image(base64.b64decode(f_info['data']), width=280)
        st.markdown(msg["content"])

# --- 6. INPUT HANDLING & VOICE TRANSCRIPTION ---
prompt_data = st.chat_input("Ask, Upload Docs/Images...", accept_file=True, accept_audio=True)

if prompt_data:
    user_text = prompt_data.text or ""
    current_files = []

    # Whisper transcription
    if prompt_data.audio:
        with st.spinner("🎤 Transcribing..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
                tmp.write(prompt_data.audio.getvalue())
                tmp_path = tmp.name
            with open(tmp_path, "rb") as af:
                tr = client.audio.transcriptions.create(file=af, model="whisper-large-v3")
                user_text = tr.text
            os.remove(tmp_path)

    # File processing loop
    if prompt_data.files:
        for f in prompt_data.files:
            if f.name.lower().endswith(('.pdf', '.pptx', '.ppt')):
                process_docs(f)
                current_files.append({"name": f.name, "type": "Document"})
            elif any(ext in f.type for ext in ["image/png", "image/jpeg"]):
                img_b64 = base64.b64encode(f.getvalue()).decode('utf-8')
                st.session_state.image_list.append({"name": f.name, "data": img_b64})
                current_files.append({"name": f.name, "type": "Image", "data": img_b64})

    if user_text or current_files:
        new_msg = {"role": "user", "content": user_text}
        if current_files: new_msg["files"] = current_files
        st.session_state.messages.append(new_msg)
        st.rerun()

# --- 7. ASSISTANT RESPONSE GENERATION ---
if st.session_state.messages and st.session_state.messages[-1]["role"] == "user":
    last_msg = st.session_state.messages[-1]
    query = last_msg["content"]

    context = ""
    if st.session_state.vector_store and query:
        qv = get_embed_model().encode([query])
        D, I = st.session_state.vector_store["index"].search(np.array(qv).astype('float32'), k=5)
        context = "\n".join([st.session_state.vector_store["chunks"][i] for i in I[0]])

    api_messages = [{"role": "system", "content": "You are Mahi's AI. Use provided context and history for your answers."}]
    for m in st.session_state.messages[-6:-1]: 
        api_messages.append({"role": m["role"], "content": m["content"]})

    if st.session_state.image_list:
        active_model = "meta-llama/llama-4-scout-17b-16e-instruct"
        content_payload = [{"type": "text", "text": f"Instruction: {query if query else 'Analyze images.'}"}]
        for img in st.session_state.image_list:
            content_payload.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img['data']}"}})
        api_messages.append({"role": "user", "content": content_payload})
    else:
        active_model = model_choice
        full_p = f"Context:\n{context}\n\nQuestion: {query}" if context else query
        api_messages.append({"role": "user", "content": full_p})

    with st.chat_message("assistant"):
        try:
            stream = client.chat.completions.create(model=active_model, messages=api_messages, stream=True)
            def parse(s):
                for c in s:
                    if c.choices[0].delta.content: yield c.choices[0].delta.content
            
            resp = st.write_stream(parse(stream))
            st.session_state.messages.append({"role": "assistant", "content": resp})
            st.session_state.image_list = [] 
            
            # Talk-Back logic
            if voice_on and resp:
                with st.spinner("🔊 Speaking..."):
                    auto_play(text_to_audio(resp, language='en'))
            
            st.rerun()
        except Exception as e:
            st.error(f"Error during generation: {e}")

# Audio re-trigger to handle browser policies
if st.session_state.last_audio:
    auto_play(st.session_state.last_audio)
    st.session_state.last_audio = None
